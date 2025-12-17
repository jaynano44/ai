from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Define a simple function to add a slide with title and content
    def add_slide(prs, layout_index, title_text, content_text=None):
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        
        if content_text:
            content = slide.placeholders[1]
            content.text = content_text
        return slide

    # 1. Title Slide
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "날씨 기반 전력량 분석 및 예측 프로젝트 합의안"
    subtitle.text = "2025년 12월 16일"

    # 2. Project Overview & Goal
    slide = add_slide(prs, 1, "1. 프로젝트 개요 및 최종 목표")
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "최종 목표: 서울시 법정동별 단기 전력 사용량 (3~7일) 예측 모델 개발 (STLF)"
    
    p = tf.add_paragraph()
    p.text = "핵심 성과:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "1. 날씨와 전력 사용량 간의 계절별 비선형 관계 분석 (히트맵)"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "2. RMSE/MAPE 기반의 모델 성능 비교"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "차별화 요소: 단일 서울 날씨 데이터에 법정동 코드를 Feature로 활용하여 지역별 패턴 분리 학습"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "현업 적용 가능성 확보"
    p.level = 1

    # 3. Data Status
    slide = add_slide(prs, 1, "2. 확보 데이터 현황 및 역할")
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "데이터 구성"
    
    data_list = [
        ("전력량", "22년 6월 ~ 25년 11월 법정동 시간별 전력량 (CSV)", "[팀원 이름] (전처리 및 병합)"),
        ("법정동 코드", "법정동 코드-명칭 매핑 자료 (XLS)", "[팀원 이름] (전처리 및 병합)"),
        ("기상 데이터", "22년 ~ 25년 서울 시간별 기상 데이터 (.HR 파일) (ZIP/CSV)", "[팀원 이름] (전처리 및 병합)")
    ]

    for name, content, role in data_list:
        p = tf.add_paragraph()
        p.text = f"{name}: {content}"
        p.level = 1
        p = tf.add_paragraph()
        p.text = f"담당: {role}"
        p.level = 2

    # 4. Roadmap Phase 1 & 2
    slide = add_slide(prs, 1, "3. 로드맵 (Phase 1 & 2)")
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    
    tf.text = "Phase 1: 데이터 통합 및 피처 엔지니어링 (D-Day 1~3)"
    p = tf.add_paragraph()
    p.text = "- 데이터 Merge: 전력량(법정동코드 문자열), 기상데이터(날짜포맷 통일)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "- 결측치 처리: 기상 데이터(NaN → 0)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "- 피처 엔지니어링: 시간 변수(요일, 월 등) 우선 생성, 추후 파생 변수(불쾌지수 등) 추가"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Phase 2: EDA 및 상관관계 분석 (D-Day 4)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "- 상관관계 히트맵: 여름(6~9월)과 겨울(11~2월) 분리 분석"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "- 지역 패턴 분석: 상업 지구(강남구) vs 주거 지구(노원구) 시간대별 패턴 비교"
    p.level = 1

    # 5. Roadmap Phase 3 & 4
    slide = add_slide(prs, 1, "3. 로드맵 (Phase 3 & 4)")
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    
    tf.text = "Phase 3: 모델링 및 성능 검증 (D-Day 5~7)"
    p = tf.add_paragraph()
    p.text = "- 훈련/테스트: 과거(22~24년) 훈련 vs 최근(25년) 검증"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "- 모델: LightGBM/XGBoost (필수) vs LSTM (선택)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "- LAG Feature: 초기 제외, 개선 단계에서 추가"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Phase 4: 결과 시뮬레이션 및 정리 (D-Day 8~10)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "- 시뮬레이션: 가상의 '다음 주' 기상 데이터로 3~7일 예측"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "- 민감도 분석: 기온 1도 상승 시 전력량 변화 분석"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "- 포트폴리오: 문제 정의 → 초기 모델 → 개선 과정 스토리텔링"
    p.level = 1

    # Save
    save_path = 'Weather_Power_Analysis_Proposal.pptx'
    prs.save(save_path)
    print(f"Presentation saved to {save_path}")

if __name__ == "__main__":
    create_presentation()
